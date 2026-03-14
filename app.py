from pptx import Presentation
from reportlab.pdfgen import canvas
from flask import send_file
import io
import pdfplumber
from flask import Flask, render_template, request
import PyPDF2
from transformers import pipeline

app = Flask(__name__)

# AI модель для конспекту
summarizer = pipeline("text-generation", model="gpt2")

def summarize(text):

    if len(text) > 500:
        text = text[:500]

    result = summarizer(text, max_length=120, num_return_sequences=1)

    generated = result[0]["generated_text"]

    # розбиваємо текст на короткі пункти
    sentences = generated.split(". ")

    bullets = []

    for s in sentences:
        s = s.strip()
        if len(s) > 5:
            bullets.append("• " + s)

    return "<br>".join(bullets)

def extract_text_from_pdf(file):
    reader = PyPDF2.PdfReader(file)
    text = ""

    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text

    return text

lecture_text = ""

@app.route("/", methods=["GET", "POST"])
def index():
    global lecture_text

    summary = ""

    if request.method == "POST":
        print("BUTTON WORKS")

        text = request.form.get("text")
        file = request.files.get("file")

        if file and file.filename.endswith(".pdf"):
            import pdfplumber

            with pdfplumber.open(file) as pdf:
                pdf_text = ""

                for page in pdf.pages:
                    pdf_text += page.extract_text()

            summary = summarize(pdf_text)
            lecture_text = pdf_text


        elif text:
            summary = summarize(text)
            lecture_text = text
        elif file and file.filename.endswith(".pptx"):

            ppt = Presentation(file)

            ppt_text = ""

            for slide in ppt.slides:
                for shape in slide.shapes:
                 if hasattr(shape, "text"):
                     ppt_text += shape.text + " "
                
            summary = summarize(ppt_text)
            lecture_text = ppt_text
    

    return render_template("index.html", summary=summary)

if __name__ == "__main__":
    app.run(debug=True)
@app.route("/download")
def download():

    summary = request.args.get("text")

    buffer = io.BytesIO()

    p = canvas.Canvas(buffer)

    y = 800

    for line in summary.split("<br>"):
        p.drawString(50, y, line)
        y -= 20

    p.save()

    buffer.seek(0)

    return send_file(
        buffer,
        as_attachment=True,
        download_name="notes.pdf",
        mimetype="application/pdf"
    )
    
@app.route("/ask", methods=["POST"])
def ask():

    question = request.form.get("question")

    context = lecture_text[:1000]

    prompt = f"Answer the question based on this lecture:\n{context}\n\nQuestion: {question}"

    result = summarizer(prompt, max_length=100, num_return_sequences=1)

    answer = result[0]["generated_text"]

    return render_template("index.html", summary="", answer=answer)